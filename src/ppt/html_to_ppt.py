"""
HTML转PPT转换器 - 使用Playwright截图并组装PPT
"""
import os
import io
import tempfile
import asyncio
from typing import List, Dict, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False
    print("Warning: Playwright not installed. HTML to PPT conversion will not work.")


class HTMLToPPTConverter:
    """
    HTML转PPT转换器
    
    核心功能:
    - 使用Playwright对HTML进行截图
    - 将截图插入到PPT幻灯片中
    - 支持批量处理多页HTML
    """
    
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    SCREENSHOT_WIDTH = 1920
    SCREENSHOT_HEIGHT = 1080
    
    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or tempfile.gettempdir()
        self._playwright = None
        self._browser = None
    
    async def convert(
        self, 
        html_slides: List[Dict[str, str]], 
        title: str = "演示文稿"
    ) -> bytes:
        """
        将HTML幻灯片列表转换为PPT
        
        Args:
            html_slides: [{"slide_index": int, "title": str, "html": str}, ...]
            title: PPT标题
            
        Returns:
            PPTX文件字节流
        """
        if not PLAYWRIGHT_AVAILABLE:
            raise RuntimeError("Playwright未安装，无法进行HTML转PPT")
        
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            
            try:
                for slide_data in html_slides:
                    html_content = slide_data.get("html", "")
                    if not html_content:
                        continue
                    
                    screenshot = await self._capture_html(browser, html_content)
                    
                    if screenshot:
                        self._add_slide_with_image(prs, screenshot)
                    else:
                        self._add_blank_slide(prs, slide_data.get("title", ""))
            
            finally:
                await browser.close()
        
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return buffer.getvalue()
    
    async def convert_single_html(
        self, 
        html_content: str, 
        output_path: str = None
    ) -> Optional[bytes]:
        """
        转换单个HTML为图片
        
        Args:
            html_content: HTML内容
            output_path: 输出路径（可选）
            
        Returns:
            图片字节流（PNG格式）
        """
        if not PLAYWRIGHT_AVAILABLE:
            raise RuntimeError("Playwright未安装")
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            
            try:
                screenshot = await self._capture_html(browser, html_content)
                
                if output_path and screenshot:
                    with open(output_path, "wb") as f:
                        f.write(screenshot)
                
                return screenshot
            
            finally:
                await browser.close()
    
    async def _capture_html(
        self, 
        browser, 
        html_content: str
    ) -> Optional[bytes]:
        """使用Playwright截取HTML页面"""
        try:
            page = await browser.new_page(
                viewport={
                    "width": self.SCREENSHOT_WIDTH,
                    "height": self.SCREENSHOT_HEIGHT
                }
            )
            
            await page.set_content(html_content, wait_until="networkidle")
            
            # HTML本身已禁用动画，直接截图
            screenshot = await page.screenshot(
                type="png",
                full_page=False
            )
            
            await page.close()
            
            return screenshot
            
        except Exception as e:
            print(f"截图失败: {e}")
            return None
    
    def _add_slide_with_image(self, prs: Presentation, image_bytes: bytes):
        """添加带有图片的幻灯片"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        image_stream = io.BytesIO(image_bytes)
        
        slide.shapes.add_picture(
            image_stream,
            Inches(0),
            Inches(0),
            width=self.SLIDE_WIDTH,
            height=self.SLIDE_HEIGHT
        )
    
    def _add_blank_slide(self, prs: Presentation, title: str):
        """添加空白幻灯片（截图失败时使用）"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        left = Inches(0.5)
        top = Inches(3)
        width = Inches(12.333)
        height = Inches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = 1
    
    def convert_sync(
        self, 
        html_slides: List[Dict[str, str]], 
        title: str = "演示文稿"
    ) -> bytes:
        """
        同步版本的转换方法
        
        Args:
            html_slides: HTML幻灯片列表
            title: PPT标题
            
        Returns:
            PPTX文件字节流
        """
        try:
            loop = asyncio.get_event_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as executor:
                future = executor.submit(
                    asyncio.run,
                    self.convert(html_slides, title)
                )
                return future.result()
        else:
            return loop.run_until_complete(self.convert(html_slides, title))
    
    def save_ppt(
        self, 
        html_slides: List[Dict[str, str]], 
        filepath: str, 
        title: str = "演示文稿"
    ) -> str:
        """
        将HTML幻灯片保存为PPT文件
        
        Args:
            html_slides: HTML幻灯片列表
            filepath: 输出文件路径
            title: PPT标题
            
        Returns:
            保存的文件路径
        """
        ppt_bytes = self.convert_sync(html_slides, title)
        
        with open(filepath, "wb") as f:
            f.write(ppt_bytes)
        
        return filepath


class HTMLSlidePreview:
    """HTML幻灯片预览器 - 用于生成预览图片"""
    
    SCREENSHOT_WIDTH = 1920
    SCREENSHOT_HEIGHT = 1080
    THUMBNAIL_WIDTH = 480
    THUMBNAIL_HEIGHT = 270
    
    def __init__(self):
        pass
    
    async def generate_preview(
        self, 
        html_content: str, 
        thumbnail: bool = True
    ) -> Dict[str, bytes]:
        """
        生成HTML预览图
        
        Args:
            html_content: HTML内容
            thumbnail: 是否生成缩略图
            
        Returns:
            {"full": bytes, "thumbnail": bytes}
        """
        if not PLAYWRIGHT_AVAILABLE:
            raise RuntimeError("Playwright未安装")
        
        result = {}
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            
            try:
                page = await browser.new_page(
                    viewport={
                        "width": self.SCREENSHOT_WIDTH,
                        "height": self.SCREENSHOT_HEIGHT
                    }
                )
                
                await page.set_content(html_content, wait_until="networkidle")
                
                # HTML本身已禁用动画，直接截图
                result["full"] = await page.screenshot(type="png")
                
                if thumbnail:
                    result["thumbnail"] = await page.screenshot(
                        type="png",
                        clip={
                            "x": 0,
                            "y": 0,
                            "width": self.SCREENSHOT_WIDTH,
                            "height": self.SCREENSHOT_HEIGHT
                        }
                    )
                
                await page.close()
            
            finally:
                await browser.close()
        
        return result
