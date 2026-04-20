# -*- coding: utf-8 -*-
"""Extract images and info from the PPT file"""
import os
import sys
from pptx import Presentation

ppt_path = r"D:\Desktop\大四学习资料\论文\PPT输出\大数据毕业论文ppt.pptx"
output_dir = r"D:\Desktop\大四学习资料\论文\PPT输出\extracted_images"
os.makedirs(output_dir, exist_ok=True)

prs = Presentation(ppt_path)

for i, slide in enumerate(prs.slides):
    print(f"\n=== Slide {i+1} ===")
    for shape in slide.shapes:
        shape_type = str(shape.shape_type)
        has_img = hasattr(shape, 'image')
        has_tf = shape.has_text_frame
        text_preview = ""
        if has_tf:
            text_preview = shape.text_frame.text[:80].replace('\n', ' ').encode('gbk', errors='replace').decode('gbk')
        
        print(f"  {shape.name.encode('gbk',errors='replace').decode('gbk')} | type={shape_type} | img={has_img} | text={text_preview}")
        
        # Extract images
        if has_img:
            ext = shape.image.content_type.split('/')[-1]
            fname = f"slide{i+1}_{shape.name.replace(' ', '_')}.{ext}"
            fpath = os.path.join(output_dir, fname)
            with open(fpath, 'wb') as f:
                f.write(shape.image.blob)
            print(f"  -> Saved: {fpath}")

print("\n\nDone!")
